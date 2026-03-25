from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─────────────────────────────────────
# SLIDE DIMENSIONS
# ─────────────────────────────────────
SW = 1920   # slide width px
SH = 1080   # slide height px

# ─────────────────────────────────────
# SCALE: px to inches
# ─────────────────────────────────────
W  = 10.0
PX = W / SW

# ─────────────────────────────────────
# HEADER / FOOTER
# ─────────────────────────────────────
HEADER_H = 162   # px
FOOTER_H = 86    # px

# ─────────────────────────────────────
# CONTENT ZONE
# ─────────────────────────────────────
MARGIN_X  = 200
CONTENT_Y = 220
CONTENT_H = SH - CONTENT_Y - FOOTER_H
TEXT_Y_OFFSET = 100  # gap between slide title and content

# ─────────────────────────────────────
# COLUMNS
# ─────────────────────────────────────
COL_GAP  = 60
COL1_X   = MARGIN_X
COL1_W   = 480
COL2_X   = COL1_X + COL1_W + COL_GAP
COL2_W   = SW - COL2_X - MARGIN_X

# ─────────────────────────────────────
# TYPOGRAPHY
# ─────────────────────────────────────
FONT_NAME    = "Inter"
FONT_TITLE   = 16
FONT_SUBHEAD = 12
FONT_BODY    = 11

# ─────────────────────────────────────
# COLOURS
# ─────────────────────────────────────
C_ACCENT = (69, 123, 157)
C_BODY   = (60, 60, 60)
C_LIGHT  = (177, 210, 218)
C_WHITE  = (255, 255, 255)
C_BG     = (255, 255, 255)


# ─────────────────────────────────────
# PRESENTATION SETUP
# ─────────────────────────────────────
def create_presentation():
    prs = Presentation()
    prs.slide_width  = Emu(9144000)
    prs.slide_height = Emu(5143500)
    return prs


# ─────────────────────────────────────
# BLANK SLIDE
# ─────────────────────────────────────
def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ─────────────────────────────────────
# BACKGROUND
# ─────────────────────────────────────
def add_bg(slide, color=C_BG):
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(W), Inches(SH * PX)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(*color)
    bg.line.fill.background()


# ─────────────────────────────────────
# HEADER / FOOTER
# ─────────────────────────────────────
def add_header_footer(slide):
    slide.shapes.add_picture(
        "PPTX_header.png",
        Inches(0), Inches(0),
        Inches(W), Inches(HEADER_H * PX)
    )
    slide.shapes.add_picture(
        "PPTX_footer.png",
        Inches(0), Inches((SH - FOOTER_H) * PX),
        Inches(W), Inches(FOOTER_H * PX)
    )


# ─────────────────────────────────────
# SOLID RECTANGLE
# ─────────────────────────────────────
def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(
        1,
        Inches(x * PX), Inches(y * PX),
        Inches(w * PX), Inches(h * PX)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.line.fill.background()


# ─────────────────────────────────────
# IMAGE
# ─────────────────────────────────────
def add_image(slide, path, x, y, w, h):
    slide.shapes.add_picture(
        path,
        Inches(x * PX), Inches(y * PX),
        Inches(w * PX), Inches(h * PX)
    )


# ─────────────────────────────────────
# TEXTBOX
# ─────────────────────────────────────
def add_textbox(slide, x, y, w, h, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(x * PX), Inches(y * PX),
        Inches(w * PX), Inches(h * PX)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    # remove internal padding
    tf.margin_left   = 0
    tf.margin_right  = 0
    tf.margin_top    = 0
    tf.margin_bottom = 0
    return tf


# ─────────────────────────────────────
# PARAGRAPH
# ─────────────────────────────────────
def add_p(tf, text, bold=False, size=None, color=None, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.font.name = FONT_NAME
    p.font.bold = bold
    p.font.size = Pt(size or FONT_BODY)
    p.font.color.rgb = RGBColor(*(color or C_BODY))
    return p


# ─────────────────────────────────────
# SLIDE TITLE
# ─────────────────────────────────────
def add_slide_title(slide, text):
    tf = add_textbox(
        slide,
        MARGIN_X, CONTENT_Y,
        SW - MARGIN_X * 2, 60
    )
    add_p(tf, text, bold=True, size=FONT_TITLE, color=C_ACCENT, first=True)