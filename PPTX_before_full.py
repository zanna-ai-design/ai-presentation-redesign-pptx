from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

# --- slide 1: title ---
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "Renewable Energy Report 2024"
subtitle.text = "GreenGrid Institute\nGlobal Capacity Analysis"

title.text_frame.paragraphs[0].font.size = Pt(40)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

# --- slide 2: chart ---
slide2 = prs.slides.add_slide(prs.slide_layouts[5])
slide2.shapes.title.text = "Renewable Energy Capacity"
slide2.shapes.add_picture(
    "chart_before.png",
    Inches(1), Inches(1.5),
    Inches(8), Inches(4.5)
)

# --- slide 3: text — typical AI generation issues ---
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
slide3.shapes.title.text = "Typical AI Generation Issues"
tf = slide3.placeholders[1].text_frame
tf.word_wrap = True

def add_heading(tf, text):
    p = tf.add_paragraph()
    p.text = text
    p.font.bold = True
    p.font.size = Pt(13)

def add_body(tf, text):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(11)

def add_gap(tf):
    p = tf.add_paragraph()
    p.text = ""

# clear default empty paragraph
tf.paragraphs[0].text = ""

add_heading(tf, "Typography and Layout")
add_body(tf, "AI-generated presentations typically rely on default fonts and slide layouts without any typographic hierarchy. Helvetica or Calibri at a uniform size is used throughout, with no distinction between headings, subheadings, and body text. Line spacing, margins, and alignment are set to defaults and rarely adjusted.")
add_gap(tf)

add_heading(tf, "Color and Visual Style")
add_body(tf, "Colors are selected arbitrarily or from built-in theme palettes that bear no relation to brand identity. Charts and diagrams often use high-saturation default colors that create visual noise rather than clarity. There is no consistent color logic across slides.")
add_gap(tf)

add_heading(tf, "Data Presentation")
add_body(tf, "Charts and tables are inserted as standalone elements without visual integration into the slide layout. Labels, legends, and axes retain their default styling. No consideration is given to how data visualization relates to the overall slide composition.")
add_gap(tf)

add_heading(tf, "Brand Identity")
add_body(tf, "AI-generated output contains no brand elements — no logo, no corporate color palette, no typography system. Each slide looks like it could belong to any organization, or none at all. The result is a document that communicates information but fails to represent the organization behind it.")

# --- slide 4: closing ---
slide4 = prs.slides.add_slide(prs.slide_layouts[0])
slide4.shapes.title.text = "Thank You"
slide4.placeholders[1].text = "GreenGrid Institute\ngreengrid.org\n© 2025"

prs.save("presentation_before.pptx")